# -*- coding: utf-8 -*-
"""Build a formal presentation for Software Practice III final review."""

from __future__ import annotations

from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT = Path(__file__).resolve().parents[1]
DOCS = PROJECT / "docs"
SCREENSHOTS = PROJECT / "screenshots" / "app"
ASSETS = DOCS / "design-assets"

OUT_MAIN = DOCS / "SmartReadAgent_期末演示PPT.pptx"
OUT_V1 = DOCS / "SmartReadAgent_V1.0_期末演示PPT.pptx"
OUT_MD = DOCS / "期末演示PPT逐页内容.md"
OUT_OUTLINE = DOCS / "PPT大纲.md"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
TEXT = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
BLUE = RGBColor(37, 99, 235)
BLUE_DARK = RGBColor(30, 64, 175)
TEAL = RGBColor(5, 150, 105)
TEAL_DARK = RGBColor(4, 120, 87)
AMBER = RGBColor(217, 119, 6)
BG = RGBColor(248, 250, 252)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(203, 213, 225)


SLIDES = [
    {
        "section": "课程项目汇报",
        "title": "SmartRead Agent",
        "subtitle": "面向学习文本的移动端智能阅读助手",
        "type": "cover",
        "notes": [
            "A03 组｜林立洲 121072021030｜江轩宇 121052023075",
            "Android / Kotlin / Jetpack Compose / LiteRT",
            "福建师范大学 2023 级软件工程 2 班",
        ],
        "image": "V1.0首页_20260512.png",
    },
    {
        "section": "背景",
        "title": "选题背景",
        "subtitle": "学生阅读资料时，最耗时间的不是打开文章，而是整理重点",
        "type": "problem",
        "cards": [
            ("阅读材料长", "教材片段、课堂资料和技术文章信息密度高，第一次阅读很难快速抓住主线。"),
            ("整理过程散", "摘要、关键词、复习点通常分散在笔记里，回看时需要重新梳理。"),
            ("移动端高频", "手机是学生最容易随手使用的设备，适合承载轻量阅读辅助流程。"),
        ],
        "image": "V1.0文本输入_20260512.png",
    },
    {
        "section": "定位",
        "title": "项目定位",
        "subtitle": "把“读文章、抓重点、问问题、做复习”串成一个闭环",
        "type": "workflow",
        "steps": [
            ("01", "输入文本", "手动粘贴或选择示例文本"),
            ("02", "智能分析", "摘要、关键词、分点摘要"),
            ("03", "端侧评分", "LiteRT 输出句子重要性"),
            ("04", "复习回看", "Agent、卡片、复习题、历史记录"),
        ],
        "image": "V1.0摘要结果_20260512.png",
    },
    {
        "section": "功能",
        "title": "V1.0 功能总览",
        "subtitle": "功能不是孤立堆叠，而是围绕阅读学习流程组织",
        "type": "matrix",
        "items": [
            ("输入", "文本输入、示例文本、字符统计"),
            ("理解", "一句话总结、关键词、分点摘要"),
            ("端侧 AI", "LiteRT 句子重要性分析"),
            ("问答", "本地规则型阅读 Agent"),
            ("复习", "知识卡片、复习题、参考答案"),
            ("回看", "历史记录保存、恢复、清空"),
        ],
        "image": "V1.0操作入口_20260512.png",
    },
    {
        "section": "界面",
        "title": "核心界面展示",
        "subtitle": "从首页到学习材料，主要页面已经形成完整操作路径",
        "type": "gallery",
        "images": [
            ("首页", "V1.0首页_20260512.png"),
            ("摘要结果", "V1.0摘要结果_20260512.png"),
            ("LiteRT 分析", "V1.0LiteRT端侧分析_20260512.png"),
            ("Agent 问答", "V1.0Agent问答_20260512.png"),
        ],
    },
    {
        "section": "架构",
        "title": "技术架构",
        "subtitle": "Android 本地分析为主，模型文件随 APK 集成",
        "type": "architecture",
        "layers": [
            ("UI 层", "Jetpack Compose", "首页、结果页、Agent、知识卡片、历史记录"),
            ("分析层", "TextAnalyzer", "摘要、关键词、分点摘要、句子切分"),
            ("智能层", "LocalReadingAgent + LiteRT", "问答意图识别、端侧句子重要性评分"),
            ("存储层", "HistoryRepository", "SharedPreferences 保存最近 12 条分析记录"),
        ],
        "image": "SmartReadAgent_模块关系图.png",
    },
    {
        "section": "端侧 AI",
        "title": "LiteRT 端侧句子重要性分析",
        "subtitle": "移动机器学习部分落实到模型文件、端侧推理和界面展示",
        "type": "ml",
        "pipeline": [
            ("文本句子", "切分文章前若干句"),
            ("5 维特征", "长度、关键词、位置、标点、提示词"),
            (".tflite 模型", "assets 集成并本地推理"),
            ("重要性等级", "分数转为高/中/低展示"),
        ],
        "image": "V1.0LiteRT端侧分析_20260512.png",
    },
    {
        "section": "问答",
        "title": "Agent 问答设计",
        "subtitle": "围绕当前文章分析结果回答，强调稳定、可解释、可复现",
        "type": "agent",
        "cards": [
            ("输入来源", "当前文章摘要、关键词、分点摘要、复习题和知识卡片。"),
            ("意图识别", "区分概括、关键词、复习重点、复习题、卡片提示等问题类型。"),
            ("回答方式", "使用本地规则组织答案，不依赖外部大模型接口。"),
        ],
        "image": "V1.0Agent快捷问题_20260512.png",
    },
    {
        "section": "复习",
        "title": "学习材料生成",
        "subtitle": "把阅读结果转化为可复习的卡片和题目",
        "type": "gallery",
        "images": [
            ("知识卡片", "V1.0知识卡片_20260512.png"),
            ("复习题", "V1.0复习题_20260512.png"),
            ("手动提问", "V1.0Agent手动提问_20260512.png"),
            ("历史恢复", "V1.0历史记录恢复_20260512.png"),
        ],
    },
    {
        "section": "数据",
        "title": "数据对象与本地存储",
        "subtitle": "分析结果、问答消息、学习材料和历史记录都有明确结构",
        "type": "architecture",
        "layers": [
            ("ArticleAnalysis", "主分析结果", "原文、摘要、关键词、分点摘要、句子重要性列表"),
            ("ChatMessage", "问答消息", "角色、内容、时间戳"),
            ("KnowledgeCard / QuizQuestion", "复习材料", "卡片标题、内容、题目、参考答案"),
            ("HistoryRecord", "历史记录", "标题、预览、保存时间、原文和摘要信息"),
        ],
        "image": "SmartReadAgent_对象模型图.png",
    },
    {
        "section": "工程",
        "title": "工程实现与版本收口",
        "subtitle": "从代码、模型、APK 到文档材料形成可复查交付",
        "type": "evidence",
        "facts": [
            ("Android 工程", "Kotlin + Compose，入口为 MainActivity"),
            ("模型文件", "sentence_importance_model.tflite 集成到 assets"),
            ("APK 产物", "SmartReadAgent-v1.0-debug.apk"),
            ("文档材料", "README、项目说明、详细设计说明书、演示脚本"),
        ],
        "image": "V1.0历史记录页_20260512.png",
    },
    {
        "section": "迭代",
        "title": "迭代过程",
        "subtitle": "从 MVP 到 V1.0，每个阶段都保留可运行成果",
        "type": "timeline",
        "milestones": [
            ("v0.1", "立项与仓库初始化"),
            ("v0.2", "文本输入与摘要 MVP"),
            ("v0.3", "Agent、知识卡片、复习题"),
            ("v0.4", "LiteRT 端侧模型集成"),
            ("v0.5", "历史记录与体验完善"),
            ("v1.0", "APK、测试、文档与汇报材料收口"),
        ],
    },
    {
        "section": "分工",
        "title": "小组分工",
        "subtitle": "双人协作完成开发、测试、演示和材料整理",
        "type": "team",
        "members": [
            ("林立洲", "项目方案、Android 核心开发、本地摘要、Agent、LiteRT 集成、版本管理、文档与汇报材料。"),
            ("江轩宇", "需求反馈、安卓真机体验、功能流程验证、UI 使用建议、文档校对和手机端操作配合。"),
        ],
    },
    {
        "section": "测试",
        "title": "测试与运行结果",
        "subtitle": "以主要流程可安装、可运行、可复现作为 V1.0 验收标准",
        "type": "review",
        "items": [
            ("功能覆盖", "摘要、关键词、分点摘要、LiteRT、Agent、知识卡片、复习题、历史记录"),
            ("运行验证", "模拟器回归验证，安卓真机完成体验反馈"),
            ("交付状态", "源码、模型、APK、截图、设计说明书和汇报材料已归档"),
        ],
        "image": "V1.0历史记录恢复_20260512.png",
    },
    {
        "section": "总结",
        "title": "总结与后续优化",
        "subtitle": "已完成课程项目闭环，同时保留可继续扩展的方向",
        "type": "summary",
        "done": [
            "完成 Android 端可运行 App",
            "形成阅读分析、端侧 AI、问答复习、历史回看的闭环",
            "完成 LiteRT 模型集成、APK 构建和材料归档",
        ],
        "next": [
            "扩大训练数据并改进句子重要性模型",
            "历史记录迁移到 Room 数据库",
            "增加 OCR 导入和可选云端增强能力",
        ],
    },
]


def set_font(run, size=18, color=TEXT, bold=False):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_box(slide, left, top, width, height, fill=CARD, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1.2)
    return box


def text_box(slide, left, top, width, height, text="", size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, color=color, bold=bold)
    return shape


def paragraph_text(shape, lines, size=17, color=TEXT, bullet=False):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
        if bullet:
            p.text = "• " + line


def add_header(slide, idx, title, subtitle, section):
    text_box(slide, Inches(0.55), Inches(0.36), Inches(1.0), Inches(0.28), f"{idx:02d}", 12, BLUE_DARK, True)
    text_box(slide, Inches(1.10), Inches(0.31), Inches(2.4), Inches(0.30), section, 10, MUTED, True)
    text_box(slide, Inches(0.55), Inches(0.75), Inches(7.8), Inches(0.54), title, 28, NAVY, True)
    text_box(slide, Inches(0.57), Inches(1.28), Inches(8.6), Inches(0.36), subtitle, 13, MUTED)
    slide.shapes.add_connector(1, Inches(0.55), Inches(1.78), Inches(12.78), Inches(1.78)).line.color.rgb = LINE


def add_footer(slide):
    text_box(slide, Inches(0.58), Inches(7.08), Inches(6.0), Inches(0.25), "SmartRead Agent｜软件项目研发实践（3）", 9, MUTED)
    text_box(slide, Inches(10.7), Inches(7.08), Inches(2.1), Inches(0.25), "A03 组 林立洲 / 江轩宇", 9, MUTED, align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, image_path, left, top, max_width, max_height, border=True):
    image_path = Path(image_path)
    if not image_path.exists():
        return None
    pic = slide.shapes.add_picture(str(image_path), left, top)
    scale = min(max_width / pic.width, max_height / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(left + (max_width - pic.width) / 2)
    pic.top = int(top + (max_height - pic.height) / 2)
    if border:
        box = add_box(slide, pic.left - Inches(0.06), pic.top - Inches(0.06), pic.width + Inches(0.12), pic.height + Inches(0.12), RGBColor(255, 255, 255), LINE)
        slide.shapes._spTree.remove(box._element)
        slide.shapes._spTree.insert(2, box._element)
    return pic


def screenshot(name: str) -> Path:
    return SCREENSHOTS / name


def asset(name: str) -> Path:
    return ASSETS / name


def render_cover(slide, data):
    text_box(slide, Inches(0.70), Inches(0.55), Inches(1.8), Inches(0.25), "课程项目汇报", 10, BLUE_DARK, True)
    text_box(slide, Inches(0.70), Inches(1.12), Inches(6.6), Inches(0.72), data["title"], 36, NAVY, True)
    text_box(slide, Inches(0.73), Inches(1.92), Inches(6.9), Inches(0.45), data["subtitle"], 18, TEXT)
    y = Inches(2.70)
    for line in data["notes"]:
        text_box(slide, Inches(0.75), y, Inches(6.6), Inches(0.32), line, 14, MUTED)
        y += Inches(0.40)
    add_box(slide, Inches(0.75), Inches(4.72), Inches(5.9), Inches(1.1), RGBColor(239, 246, 255), RGBColor(147, 197, 253))
    paragraph_text(slide.shapes[-1], [
        "项目目标：在 Android 端完成文本阅读分析、端侧句子重要性评分、问答复习和历史回看。",
        "汇报重点：应用价值、功能完整度、LiteRT 集成、工程实现和团队协作。"
    ], 13, BLUE_DARK)
    add_picture_fit(slide, screenshot(data["image"]), Inches(8.2), Inches(0.58), Inches(3.0), Inches(6.4))


def render_problem(slide, data):
    x = Inches(0.75)
    for title, body in data["cards"]:
        add_box(slide, x, Inches(2.10), Inches(2.75), Inches(1.50), RGBColor(239, 246, 255), RGBColor(147, 197, 253))
        text_box(slide, x + Inches(0.22), Inches(2.30), Inches(2.3), Inches(0.32), title, 18, BLUE_DARK, True)
        text_box(slide, x + Inches(0.22), Inches(2.78), Inches(2.3), Inches(0.58), body, 12.5, TEXT)
        x += Inches(3.05)
    add_box(slide, Inches(0.78), Inches(4.10), Inches(5.95), Inches(1.28), RGBColor(240, 253, 250), RGBColor(94, 234, 212))
    paragraph_text(slide.shapes[-1], ["设计取向：不做泛泛的聊天工具，而是聚焦“课程资料阅读”这一高频小场景。", "V1.0 的价值在于把阅读、理解和复习整理放进一个移动端闭环。"], 14, TEAL_DARK)
    add_picture_fit(slide, screenshot(data["image"]), Inches(8.35), Inches(2.05), Inches(2.6), Inches(4.4))


def render_workflow(slide, data):
    x = Inches(0.75)
    colors = [RGBColor(239, 246, 255), RGBColor(236, 253, 245), RGBColor(255, 251, 235), RGBColor(240, 249, 255)]
    for idx, (no, title, body) in enumerate(data["steps"]):
        add_box(slide, x, Inches(2.10), Inches(2.55), Inches(1.55), colors[idx], LINE)
        text_box(slide, x + Inches(0.20), Inches(2.28), Inches(0.55), Inches(0.30), no, 18, BLUE_DARK, True)
        text_box(slide, x + Inches(0.82), Inches(2.25), Inches(1.6), Inches(0.32), title, 17, NAVY, True)
        text_box(slide, x + Inches(0.25), Inches(2.84), Inches(2.1), Inches(0.42), body, 12.5, MUTED)
        if idx < 3:
            text_box(slide, x + Inches(2.48), Inches(2.70), Inches(0.28), Inches(0.24), "→", 18, MUTED, True)
        x += Inches(2.82)
    add_picture_fit(slide, screenshot(data["image"]), Inches(4.95), Inches(4.15), Inches(3.7), Inches(2.2))


def render_matrix(slide, data):
    positions = [(0.75, 2.10), (3.75, 2.10), (6.75, 2.10), (0.75, 4.25), (3.75, 4.25), (6.75, 4.25)]
    for (title, body), (x, y) in zip(data["items"], positions):
        add_box(slide, Inches(x), Inches(y), Inches(2.62), Inches(1.36), RGBColor(255, 255, 255), LINE)
        text_box(slide, Inches(x + 0.22), Inches(y + 0.18), Inches(2.1), Inches(0.30), title, 17, TEAL_DARK, True)
        text_box(slide, Inches(x + 0.22), Inches(y + 0.60), Inches(2.1), Inches(0.48), body, 12.5, TEXT)
    add_picture_fit(slide, screenshot(data["image"]), Inches(9.75), Inches(2.00), Inches(2.35), Inches(4.80))


def render_gallery(slide, data):
    positions = [(0.80, 2.05), (3.98, 2.05), (7.16, 2.05), (10.34, 2.05)]
    for (label, img), (x, y) in zip(data["images"], positions):
        add_box(slide, Inches(x), Inches(y), Inches(2.38), Inches(4.55), RGBColor(255, 255, 255), LINE)
        add_picture_fit(slide, screenshot(img), Inches(x + 0.18), Inches(y + 0.18), Inches(2.02), Inches(3.72), border=False)
        text_box(slide, Inches(x + 0.22), Inches(y + 4.05), Inches(1.95), Inches(0.30), label, 13, BLUE_DARK, True, PP_ALIGN.CENTER)


def render_architecture(slide, data):
    y = Inches(2.05)
    for title, tech, body in data["layers"]:
        add_box(slide, Inches(0.78), y, Inches(5.25), Inches(0.82), RGBColor(255, 255, 255), LINE)
        text_box(slide, Inches(1.02), y + Inches(0.15), Inches(1.2), Inches(0.25), title, 14, BLUE_DARK, True)
        text_box(slide, Inches(2.18), y + Inches(0.15), Inches(1.75), Inches(0.25), tech, 13, NAVY, True)
        text_box(slide, Inches(4.0), y + Inches(0.15), Inches(1.75), Inches(0.28), body, 10.5, MUTED)
        y += Inches(0.98)
    img = asset(data["image"]) if data["image"].endswith(".png") else screenshot(data["image"])
    add_picture_fit(slide, img, Inches(6.55), Inches(2.0), Inches(5.9), Inches(4.45))


def render_ml(slide, data):
    x = Inches(0.78)
    for title, body in data["pipeline"]:
        add_box(slide, x, Inches(2.20), Inches(2.45), Inches(1.20), RGBColor(239, 246, 255), RGBColor(147, 197, 253))
        text_box(slide, x + Inches(0.20), Inches(2.38), Inches(2.0), Inches(0.30), title, 15, BLUE_DARK, True, PP_ALIGN.CENTER)
        text_box(slide, x + Inches(0.22), Inches(2.82), Inches(1.95), Inches(0.28), body, 10.5, TEXT, align=PP_ALIGN.CENTER)
        x += Inches(2.72)
    text_box(slide, Inches(1.65), Inches(3.55), Inches(9.1), Inches(0.40), "Text → Feature Vector → LiteRT Interpreter → Score / Level", 18, NAVY, True, PP_ALIGN.CENTER)
    add_picture_fit(slide, screenshot(data["image"]), Inches(4.7), Inches(4.05), Inches(3.7), Inches(2.35))


def render_agent(slide, data):
    x = Inches(0.85)
    for title, body in data["cards"]:
        add_box(slide, x, Inches(2.10), Inches(3.05), Inches(1.45), RGBColor(255, 255, 255), LINE)
        text_box(slide, x + Inches(0.20), Inches(2.30), Inches(2.4), Inches(0.30), title, 16, TEAL_DARK, True)
        text_box(slide, x + Inches(0.20), Inches(2.78), Inches(2.55), Inches(0.48), body, 12, TEXT)
        x += Inches(3.35)
    add_picture_fit(slide, screenshot(data["image"]), Inches(4.8), Inches(4.0), Inches(3.8), Inches(2.45))


def render_evidence(slide, data):
    y = Inches(2.05)
    for title, body in data["facts"]:
        add_box(slide, Inches(0.78), y, Inches(5.6), Inches(0.82), RGBColor(255, 255, 255), LINE)
        text_box(slide, Inches(1.03), y + Inches(0.17), Inches(1.45), Inches(0.28), title, 14, BLUE_DARK, True)
        text_box(slide, Inches(2.62), y + Inches(0.17), Inches(3.2), Inches(0.28), body, 12, TEXT)
        y += Inches(0.97)
    add_picture_fit(slide, screenshot(data["image"]), Inches(8.1), Inches(2.03), Inches(2.75), Inches(4.5))


def render_timeline(slide, data):
    x = Inches(0.80)
    y = Inches(3.00)
    slide.shapes.add_connector(1, Inches(1.05), y + Inches(0.35), Inches(12.1), y + Inches(0.35)).line.color.rgb = LINE
    for version, body in data["milestones"]:
        add_box(slide, x, y, Inches(1.75), Inches(1.15), RGBColor(255, 255, 255), LINE)
        text_box(slide, x + Inches(0.12), y + Inches(0.16), Inches(1.45), Inches(0.28), version, 15, BLUE_DARK, True, PP_ALIGN.CENTER)
        text_box(slide, x + Inches(0.12), y + Inches(0.58), Inches(1.45), Inches(0.36), body, 10.8, TEXT, align=PP_ALIGN.CENTER)
        x += Inches(2.02)
    add_box(slide, Inches(1.2), Inches(5.15), Inches(10.9), Inches(0.95), RGBColor(240, 253, 250), RGBColor(94, 234, 212))
    paragraph_text(slide.shapes[-1], ["迭代策略：每个阶段先保证可运行，再补充测试截图、版本记录和提交材料。"], 15, TEAL_DARK)


def render_team(slide, data):
    x = Inches(0.85)
    for name, body in data["members"]:
        add_box(slide, x, Inches(2.15), Inches(5.35), Inches(2.25), RGBColor(255, 255, 255), LINE)
        text_box(slide, x + Inches(0.30), Inches(2.45), Inches(4.7), Inches(0.35), name, 22, NAVY, True)
        text_box(slide, x + Inches(0.30), Inches(3.08), Inches(4.65), Inches(0.85), body, 14, TEXT)
        x += Inches(5.8)
    add_box(slide, Inches(2.2), Inches(5.15), Inches(8.95), Inches(0.92), RGBColor(239, 246, 255), RGBColor(147, 197, 253))
    paragraph_text(slide.shapes[-1], ["协作方式：开发、真机体验、文档校对和演示操作分开推进，最终统一到 V1.0 交付材料。"], 15, BLUE_DARK)


def render_review(slide, data):
    y = Inches(2.05)
    for title, body in data["items"]:
        add_box(slide, Inches(0.82), y, Inches(5.9), Inches(1.05), RGBColor(255, 255, 255), LINE)
        text_box(slide, Inches(1.05), y + Inches(0.20), Inches(1.4), Inches(0.28), title, 15, TEAL_DARK, True)
        text_box(slide, Inches(2.42), y + Inches(0.20), Inches(3.85), Inches(0.45), body, 12, TEXT)
        y += Inches(1.22)
    add_picture_fit(slide, screenshot(data["image"]), Inches(8.25), Inches(2.0), Inches(2.7), Inches(4.65))


def render_summary(slide, data):
    add_box(slide, Inches(0.82), Inches(2.10), Inches(5.75), Inches(3.1), RGBColor(240, 253, 250), RGBColor(94, 234, 212))
    text_box(slide, Inches(1.12), Inches(2.40), Inches(4.7), Inches(0.35), "已经完成", 20, TEAL_DARK, True)
    paragraph_text(text_box(slide, Inches(1.12), Inches(2.95), Inches(4.85), Inches(1.65)), data["done"], 14, TEXT, True)
    add_box(slide, Inches(7.00), Inches(2.10), Inches(5.25), Inches(3.1), RGBColor(255, 251, 235), RGBColor(252, 211, 77))
    text_box(slide, Inches(7.30), Inches(2.40), Inches(4.3), Inches(0.35), "后续方向", 20, AMBER, True)
    paragraph_text(text_box(slide, Inches(7.30), Inches(2.95), Inches(4.35), Inches(1.65)), data["next"], 14, TEXT, True)
    text_box(slide, Inches(2.0), Inches(5.85), Inches(9.4), Inches(0.5), "SmartRead Agent V1.0 已形成可运行、可演示、可复查的课程项目交付。", 20, NAVY, True, PP_ALIGN.CENTER)


def render_slide(prs, idx, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    if data["type"] != "cover":
        add_header(slide, idx, data["title"], data["subtitle"], data["section"])
        add_footer(slide)
    kind = data["type"]
    if kind == "cover":
        render_cover(slide, data)
    elif kind == "problem":
        render_problem(slide, data)
    elif kind == "workflow":
        render_workflow(slide, data)
    elif kind == "matrix":
        render_matrix(slide, data)
    elif kind == "gallery":
        render_gallery(slide, data)
    elif kind == "architecture":
        render_architecture(slide, data)
    elif kind == "ml":
        render_ml(slide, data)
    elif kind == "agent":
        render_agent(slide, data)
    elif kind == "evidence":
        render_evidence(slide, data)
    elif kind == "timeline":
        render_timeline(slide, data)
    elif kind == "team":
        render_team(slide, data)
    elif kind == "review":
        render_review(slide, data)
    elif kind == "summary":
        render_summary(slide, data)
    return slide


def build_markdown() -> str:
    lines = ["# SmartRead Agent 正式汇报 PPT 逐页内容", ""]
    for idx, slide in enumerate(SLIDES, 1):
        lines.extend([f"## {idx}. {slide['title']}", "", f"副标题：{slide['subtitle']}", ""])
        if "cards" in slide:
            for title, body in slide["cards"]:
                lines.append(f"- {title}：{body}")
        if "steps" in slide:
            for no, title, body in slide["steps"]:
                lines.append(f"- {no} {title}：{body}")
        if "items" in slide:
            for title, body in slide["items"]:
                lines.append(f"- {title}：{body}")
        if "layers" in slide:
            for title, tech, body in slide["layers"]:
                lines.append(f"- {title} / {tech}：{body}")
        if "pipeline" in slide:
            for title, body in slide["pipeline"]:
                lines.append(f"- {title}：{body}")
        if "facts" in slide:
            for title, body in slide["facts"]:
                lines.append(f"- {title}：{body}")
        if "milestones" in slide:
            for title, body in slide["milestones"]:
                lines.append(f"- {title}：{body}")
        if "members" in slide:
            for title, body in slide["members"]:
                lines.append(f"- {title}：{body}")
        if "done" in slide:
            lines.append("- 已完成：" + "；".join(slide["done"]))
            lines.append("- 后续方向：" + "；".join(slide["next"]))
        lines.append("")
    return "\n".join(lines)


def main() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for idx, data in enumerate(SLIDES, 1):
        render_slide(prs, idx, data)
    prs.save(OUT_MAIN)
    shutil.copy2(OUT_MAIN, OUT_V1)
    markdown = build_markdown()
    OUT_MD.write_text(markdown, encoding="utf-8")
    OUT_OUTLINE.write_text(markdown, encoding="utf-8")
    print(OUT_MAIN)
    print(OUT_V1)
    print(OUT_MD)
    print(OUT_OUTLINE)


if __name__ == "__main__":
    main()
