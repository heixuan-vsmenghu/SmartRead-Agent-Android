# -*- coding: utf-8 -*-
"""Generate final Software Practice III presentation and design materials.

The script keeps the deliverables reproducible: markdown files are the source of
truth, while PPTX/DOCX are generated from the same curated content.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path
import re
import subprocess
import sys

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


PROJECT = Path(__file__).resolve().parents[1]
ROOT = PROJECT.parent
DOCS = PROJECT / "docs"
DEMO = PROJECT / "demo"
SCREENSHOTS = PROJECT / "screenshots" / "app"
RELEASE = PROJECT / "release"
PM_MEMO = (
    PROJECT.parent.parent
    / "软件项目管理"
    / "期末大作业_SmartReadAgent_项目管理"
    / "08_个人项目管理报告"
    / "个人项目管理报告_素材备忘.md"
)


TODAY = date(2026, 6, 29).isoformat()


SLIDES = [
    {
        "title": "SmartRead Agent：AI 智读助手",
        "subtitle": "第4组 / Android Kotlin 移动端 AI 阅读辅助应用",
        "bullets": [
            "课程：软件项目研发实践（3）",
            "演示时间：2026-07-01 9:50-10:05，计网大楼 202",
            "成员：林立洲（汇报、开发、文档），江轩宇（演示、测试、反馈）",
        ],
        "image": "V1.0首页_20260512.png",
    },
    {
        "title": "项目背景：学生长文本阅读成本高",
        "subtitle": "从课堂资料、教材片段、技术文章中快速抓重点",
        "bullets": [
            "课程学习中常见资料篇幅长，手工整理摘要和复习点耗时",
            "手机端场景适合做轻量、随手可用的阅读辅助工具",
            "目标不是替代学习，而是降低第一遍理解和复习整理成本",
        ],
        "image": "V1.0文本输入_20260512.png",
    },
    {
        "title": "项目目标：形成阅读分析闭环",
        "subtitle": "输入文本 -> 智能分析 -> 问答复习 -> 历史恢复",
        "bullets": [
            "支持文本输入和示例文本，方便现场稳定演示",
            "生成一句话总结、关键词和分点摘要",
            "接入 LiteRT 端侧句子重要性分析，突出移动机器学习",
            "提供 Agent 问答、知识卡片、复习题和历史记录",
        ],
        "image": "V1.0摘要结果_20260512.png",
    },
    {
        "title": "功能总览",
        "subtitle": "围绕阅读、理解、复习三类动作组织功能",
        "bullets": [
            "阅读：示例文本、手动输入、摘要结果展示",
            "理解：关键词、分点摘要、LiteRT 句子重要性评分",
            "复习：Agent 问答、知识卡片、复习题生成",
            "回看：本地历史记录保存、恢复和清空",
        ],
        "image": "V1.0操作入口_20260512.png",
    },
    {
        "title": "技术架构",
        "subtitle": "Android 本地完成主要分析流程",
        "bullets": [
            "Android/Kotlin + Jetpack Compose 构建移动端界面",
            "TextAnalyzer 完成本地摘要、关键词和分点摘要",
            "LocalReadingAgent 作为本地规则型阅读 Agent",
            "SentenceImportanceClassifier 加载 assets 中的 .tflite 模型",
            "SharedPreferences + JSON 保存最近 12 条历史记录",
        ],
        "diagram": True,
    },
    {
        "title": "机器学习整合：LiteRT 端侧分析",
        "subtitle": "评分表中机器学习整合占 40%，这里是重点",
        "bullets": [
            "模型任务：判断句子在当前文章中的重要程度",
            "输入 5 维特征：长度、关键词重合、位置、标点、提示词",
            "模型结构：Input(5) -> Dense(8) -> Dense(4) -> Dense(1)",
            ".tflite 文件放入 Android assets，推理不依赖网络",
            "模型加载失败时使用规则评分兜底，保证现场演示稳定",
        ],
        "image": "V1.0LiteRT端侧分析_20260512.png",
    },
    {
        "title": "Agent 问答：本地规则型阅读 Agent",
        "subtitle": "当前版本没有调用 GPT API，也没有依赖云端大模型",
        "bullets": [
            "Agent 基于当前文章的摘要、关键词、分点摘要和复习题规则回答",
            "通过问题意图识别区分概括、关键词、复习重点、复习题和知识卡片",
            "优点是离线可用、稳定、可解释，适合课程现场演示",
            "后续可把云端大模型作为增强能力，而不是当前版本依赖项",
        ],
        "image": "V1.0Agent问答_20260512.png",
    },
    {
        "title": "现场 App 演示流程",
        "subtitle": "江轩宇负责手机操作，林立洲同步说明",
        "bullets": [
            "打开 App，点击示例文本，执行开始智能分析",
            "展示摘要、关键词、分点摘要和 LiteRT 端侧评分",
            "进入 Agent 问答页，点击快捷问题并展示回答",
            "进入知识卡片页，展示知识卡片和复习题",
            "进入历史记录页，恢复一次历史分析后回到首页",
        ],
        "image": "V1.0Agent快捷问题_20260512.png",
    },
    {
        "title": "项目迭代过程",
        "subtitle": "从 MVP 到 V1.0 演示稳定版",
        "bullets": [
            "v0.1：立项、目录、Git/Gitee/GitHub 初始化",
            "v0.2：文本输入、示例文本、本地摘要 MVP",
            "v0.3：Agent 问答、知识卡片、复习题",
            "v0.4：LiteRT 端侧模型集成与分析展示",
            "v0.5-v1.0：历史记录、真机体验、APK 和演示材料收口",
        ],
        "image": "V1.0历史记录页_20260512.png",
    },
    {
        "title": "团队协作与分工",
        "subtitle": "按真实参与情况说明，不夸大代码贡献",
        "bullets": [
            "林立洲：选题方案、Android 核心功能、Agent、LiteRT、版本管理和材料整理",
            "江轩宇：需求反馈、真机安装、功能体验测试、UI 建议、文档校对和现场操作",
            "使用 Gitee/GitHub 记录版本迭代，使用 QQ 和项目管理平台保留过程证据",
            "演示中林立洲主讲，江轩宇配合操作，答辩由林立洲主答",
        ],
        "image": None,
    },
    {
        "title": "测试与运行结果",
        "subtitle": "以能安装、能运行、能演示为 V1.0 收口标准",
        "bullets": [
            "完成模拟器测试：摘要、LiteRT、Agent、知识卡片、历史记录",
            "生成课程演示用 debug APK：SmartReadAgent-v1.0-debug.apk",
            "完成安卓真机体验反馈，确认主要流程可进入",
            "当前暂无阻塞性 Bug，后续优化集中在模型效果和存储结构",
        ],
        "image": "V1.0历史记录恢复_20260512.png",
    },
    {
        "title": "项目亮点",
        "subtitle": "对照评分项突出应用价值和机器学习整合",
        "bullets": [
            "场景真实：面向学生阅读教材、课堂资料和技术文章",
            "功能闭环：摘要、问答、卡片、复习题、历史记录形成学习流程",
            "移动机器学习：LiteRT 模型集成到 Android 端侧执行",
            "演示稳定：不依赖网络和云端 API，课堂环境可控",
            "过程完整：版本、截图、测试、APK、文档材料均已归档",
        ],
        "image": "V1.0知识卡片_20260512.png",
    },
    {
        "title": "不足与后续优化",
        "subtitle": "客观说明边界，不虚构未实现功能",
        "bullets": [
            "摘要和 Agent 当前偏轻量，主要依赖本地规则和当前文章分析结果",
            "LiteRT 模型规模较小，适合课程演示，准确率仍可提升",
            "历史记录当前使用 SharedPreferences，后续可迁移到 Room 数据库",
            "后续可增加 OCR、正式签名 release 包、云端大模型增强能力",
        ],
        "image": None,
    },
    {
        "title": "谢谢观看",
        "subtitle": "SmartRead Agent V1.0 期末演示稳定版",
        "bullets": [
            "GitHub/Gitee：用于源码、版本和 Release 材料展示",
            "APK：课程演示用 debug APK，不作为商业发布版",
            "详细设计说明书：电子版和纸质版 2026-07-10 前提交",
            "欢迎老师和同学提问",
        ],
        "image": "V1.0首页_20260512.png",
    },
]


QA = [
    ("这个项目主要解决什么问题？", "它解决的是学生阅读教材、课堂资料和技术文章时整理重点效率低的问题。用户把文本放进 App 后，可以快速得到一句话总结、关键词、分点摘要，再通过 Agent 问答、知识卡片和复习题继续复习。"),
    ("为什么选择阅读摘要这个方向？", "这个方向和学生自己的学习场景贴得比较近，演示时也容易验证。相比做一个很大的泛 AI 应用，阅读摘要更适合在移动端做出完整闭环，也能自然接入端侧模型。"),
    ("为什么叫 SmartRead Agent？", "SmartRead 表示智能阅读，Agent 表示它不是只显示摘要，还会根据当前文章分析结果响应用户问题，辅助整理复习重点。当前版本是轻量的本地规则型阅读 Agent。"),
    ("Agent 是否调用 GPT API？", "没有。当前版本没有调用 GPT API，也没有依赖云端大模型。Agent 问答基于当前文章的一句话总结、关键词、分点摘要和复习题生成规则，通过问题意图识别返回回答。"),
    ("Agent 具体如何实现？", "代码中由 LocalReadingAgent 负责。它先判断用户问题属于概括、关键词、复习重点、复习题还是知识卡片等意图，然后调用已有的分析结果组织回答。它的优势是离线可用、稳定、可解释。"),
    ("LiteRT 模型做了什么？", "LiteRT 模型用于句子重要性分析。App 对文章前几句提取 5 维特征，送入 .tflite 模型，输出 0 到 1 的重要性分数，并在界面上显示高、中、低等级和推理来源。"),
    ("为什么不用云端大模型？", "课程演示更重视移动端机器学习落地和现场稳定性。云端大模型效果更强，但需要网络、账号和 API Key，不利于课堂环境稳定演示。当前版本先实现本地能力，后续可以作为增强项接入云端模型。"),
    ("端侧推理有什么优势？", "端侧推理不依赖网络，响应稳定，也能避免把文章内容上传到服务器。对课程项目来说，它能清楚展示模型文件随 APK 集成并在手机端运行。"),
    ("模型输入特征是什么？", "模型输入是 5 维浮点特征：句子长度归一化、关键词重合度、句子位置分数、标点提示分数、总结提示词分数。"),
    (".tflite 文件如何集成到 Android？", "模型文件 sentence_importance_model.tflite 放在 android/app/src/main/assets 目录。SentenceImportanceClassifier 通过 context.assets.open 读取模型，再用 TensorFlow Lite Interpreter 执行推理。"),
    ("LiteRT 和普通规则算法有什么区别？", "普通规则算法是人工设定权重直接打分；LiteRT 是把特征输入到训练后导出的模型，由模型输出重要性分数。当前项目还保留规则兜底，保证模型异常时 App 不崩溃。"),
    ("如果模型加载失败怎么办？", "SentenceImportanceClassifier 会检测 Interpreter 是否可用。如果模型加载或推理失败，就使用 fallbackScore 规则评分，界面会显示规则评分兜底，保证演示流程继续。"),
    ("摘要算法如何实现？", "TextAnalyzer 先切分句子，再根据关键词命中、句子长度和位置给句子打分，选出一句话总结，并挑选得分较高的句子作为分点摘要。"),
    ("关键词如何提取？", "关键词提取结合领域词、英文技术词和中文短语窗口。代码会过滤部分停用词，再按得分排序，并避免关键词之间过度重复。"),
    ("历史记录如何保存？", "HistoryRepository 使用 SharedPreferences 保存 JSON 字符串，最多保留最近 12 条分析记录，包括原文、摘要、关键词、字符数、句子数和保存时间。"),
    ("为什么没有使用数据库？", "V1.0 主要面向课程演示，历史记录数据量很小，SharedPreferences 足够支撑最近记录保存和恢复。后续如果要做正式版本，可以迁移到 Room 数据库。"),
    ("App 是否支持离线？", "支持。摘要、关键词、Agent 问答、知识卡片、复习题、历史记录和 LiteRT 端侧推理都在本地完成，不需要网络。"),
    ("你们怎么测试的？", "主要做了模拟器测试和安卓真机体验测试。测试覆盖 App 启动、示例文本分析、摘要结果、LiteRT 区域、Agent 问答、知识卡片、复习题、历史记录恢复和清空。"),
    ("小组怎么分工？", "林立洲负责项目方案、Android 核心开发、本地摘要、Agent、LiteRT 集成、版本管理和材料整理。江轩宇负责需求反馈、真机安装与功能验证、UI 使用建议、文档校对和现场操作演示。"),
    ("江轩宇具体参与了什么？", "江轩宇参与需求反馈、功能体验测试、安卓真机安装验证、Agent 问答、知识卡片和历史记录等流程的体验确认，也负责演示时手机端操作配合。"),
    ("项目管理平台怎么用？", "我们用项目管理平台记录阶段任务、任务状态、Bug 和进度复盘，配合 QQ 沟通截图和测试记录，形成过程证据。"),
    ("Gitee 提交记录体现了什么？", "提交记录体现了从 v0.1 立项、v0.2 摘要 MVP、v0.3 Agent、v0.4 LiteRT、v0.5 历史记录到 v1.0 稳定版的迭代过程。"),
    ("项目创新点是什么？", "创新点在于把阅读摘要、复习材料生成和 LiteRT 端侧句子重要性分析结合在一个移动端 App 中，场景聚焦学生学习，不是只做单一摘要 Demo。"),
    ("项目不足是什么？", "摘要和 Agent 当前比较轻量，模型训练数据也偏课程演示规模。它能稳定展示端侧分析流程，但距离商业级智能阅读产品还有差距。"),
    ("后续如何优化？", "后续可以提升模型数据和训练方式，引入 Room 数据库，增加 OCR 导入，完善正式签名 release 包，并在用户允许的情况下接入云端大模型增强问答效果。"),
    ("为什么机器学习部分没有做很大的模型？", "课程项目周期和移动端资源有限，选择轻量模型更适合落地到 APK。我们重点展示从特征设计、模型导出、assets 集成到端侧推理展示的完整链路。"),
    ("模型输出分数怎么解释？", "分数越接近 1，表示句子越可能是当前文章中的重要句子。界面把分数转换为高、中、低等级，方便演示和普通用户理解。"),
    ("演示时如果手机出问题怎么办？", "准备了模拟器和截图备份，也保留了 APK 和演示脚本。现场如果真机不稳定，可以切换到模拟器或用截图说明关键流程。"),
    ("为什么说它适合课程演示？", "因为核心能力本地运行，流程短，功能闭环清楚，机器学习部分有实际模型文件和端侧推理展示，同时也能体现项目迭代和小组协作过程。"),
    ("详细设计说明书准备提交什么？", "提交 SmartReadAgent_软件详细设计说明书_初稿.docx 和对应 Markdown，内容包括架构、模块、算法、数据结构、界面、测试、部署和后续优化。"),
]


DESIGN_MD = f"""# SmartRead Agent 软件详细设计说明书（初稿）

## 文档信息

| 项目 | 内容 |
|---|---|
| 项目名称 | SmartRead Agent：AI 智读助手 |
| 文档版本 | v1.0 初稿 |
| 编写日期 | {TODAY} |
| 课程名称 | 软件项目研发实践（3） |
| 指导老师 | 林立老师 |
| 演示顺序 | 第4组，2026-07-01 9:50-10:05，计网大楼 202 |
| 小组成员 | 林立洲（121072021030）、江轩宇（121052023075） |
| 适用范围 | 课程期末演示、详细设计说明书提交、项目复盘 |

## 修订记录

| 版本 | 日期 | 修订内容 | 修订人 |
|---|---|---|---|
| v0.1 | 2026-05-11 | 项目立项、目录和版本管理初始化 | 林立洲 |
| v0.2 | 2026-05-12 | 完成文本输入、本地摘要和关键词 MVP | 林立洲 |
| v0.3 | 2026-05-12 | 增加 Agent 问答、知识卡片和复习题 | 林立洲 |
| v0.4 | 2026-05-12 | 集成 LiteRT 端侧句子重要性分析模型 | 林立洲 |
| v0.5 | 2026-05-12 | 增加本地历史记录保存、恢复和清空 | 林立洲 |
| v1.0 | 2026-05-13 | 收口期末演示稳定版、APK、截图和说明材料 | 林立洲、江轩宇 |

## 目录

1. 引言
2. 项目概述
3. 总体设计
4. 模块详细设计
5. 核心算法与模型设计
6. 数据结构设计
7. 界面设计
8. 测试设计
9. 部署与运行
10. 项目总结与后续优化

## 1 引言

### 1.1 编写目的

本说明书用于记录 SmartRead Agent V1.0 期末演示稳定版的详细设计。文档面向《软件项目研发实践（3）》课程期末考核，重点说明系统的功能范围、模块划分、数据结构、端侧机器学习模型、界面流程、测试情况和部署方式。它不是商业产品发布文档，而是课程项目在完成移动端 AI 应用开发后，对设计和实现过程的一次整理。

本项目已经完成课程演示用 debug APK，文件位于 `release/SmartReadAgent-v1.0-debug.apk`。该 APK 用于课堂验收、真机测试和现场演示，不用于应用商店正式商业发布。文档中涉及的 OCR、云端大模型增强、Room 数据库和正式签名 release 包均作为后续优化方向，不写成当前已经实现的功能。

### 1.2 项目背景

在日常学习中，学生经常需要阅读教材段落、课堂讲义、技术文章和项目资料。资料内容往往较长，手工划重点、写摘要、整理复习题会消耗较多时间。SmartRead Agent 的设计出发点是做一个轻量的移动端阅读辅助工具，让用户在手机上粘贴或选择示例文本后，快速得到摘要、关键词、分点摘要和复习材料。

课程评分中机器学习在应用中的整合程度占比较高，因此项目没有停留在简单文本处理 Demo，而是在 V0.4 阶段加入 LiteRT / TensorFlow Lite 端侧句子重要性分析。模型文件随 APK 放入 Android assets，运行时在本地加载并推理，体现移动端机器学习应用的完整链路。

### 1.3 读者对象

本文档主要面向课程指导老师、参与评分的班委和组长、小组成员，以及后续查看项目的同学。读者不需要提前了解完整代码，但应能通过本文档理解系统由哪些模块组成、各模块如何协作、机器学习部分如何接入 Android 应用。

### 1.4 参考资料

- Android 官方文档与 Kotlin 语言基础资料。
- Jetpack Compose UI 开发资料。
- TensorFlow Lite / LiteRT Android 推理相关资料。
- 项目仓库中的 `README.md`、`docs/技术方案.md`、`docs/V0.4_LiteRT端侧模型集成说明.md`。
- 课程文件：`2023级软工期末考核说明.pdf`、`2023级软工期末课程演示流程.pdf`、`2023级期末演示评分表.xlsx`。

### 1.5 术语说明

- LiteRT / TensorFlow Lite：用于移动端和端侧设备的轻量模型推理框架。
- 本地规则型阅读 Agent：本项目当前版本的 Agent 形式，不调用 GPT API，不依赖云端大模型，而是根据当前文章的摘要、关键词、分点摘要和复习题规则进行问答。
- 端侧推理：模型在手机本地运行，不把文本上传到服务器。
- APK：Android 安装包，本项目 V1.0 使用 debug APK 进行课程演示。

## 2 项目概述

### 2.1 系统目标

SmartRead Agent 的目标是完成一个可运行、可演示、能体现移动端机器学习能力的 Android 应用。系统需要支持文本输入与示例文本、自动摘要、关键词提取、分点摘要、LiteRT 句子重要性分析、Agent 问答、知识卡片、复习题和历史记录。它的核心价值不是做一个大而全的学习平台，而是在短时间内帮助用户完成“读文章、抓重点、问问题、复习回看”的闭环。

### 2.2 应用场景

典型场景包括：学生阅读课程讲义时先用 App 得到摘要；阅读技术文章时快速提取关键词和重点句；复习前用知识卡片和复习题检查理解；课堂演示时使用示例文本稳定展示各功能。由于当前能力主要在本地运行，现场演示时不依赖网络，能够降低环境不确定性。

### 2.3 功能范围

V1.0 已实现的功能包括：文本输入、示例文本、一句话总结、关键词、分点摘要、LiteRT 端侧句子重要性分析、本地规则型阅读 Agent 问答、知识卡片、复习题、本地历史记录保存、历史记录恢复和清空历史记录。

V1.0 未实现的功能包括：用户登录、云端同步、OCR 图片识别、语音输入、多设备同步、正式商业发布包、服务端接口和云端大模型 API。这些内容不在当前演示范围内，后续可以作为扩展方向。

### 2.4 非功能需求

系统需要满足稳定演示、离线可用、操作路径短、界面信息清晰、模型加载失败可兜底、APK 可安装和主要流程可复现等要求。由于演示时间只有 15 分钟，功能入口必须集中，演示文本必须提前准备，手机端操作不能依赖复杂配置。

## 3 总体设计

### 3.1 系统总体架构

系统采用单体 Android 应用结构，核心代码集中在 `android/app/src/main/java/com/smartread/agent/`。界面层使用 Jetpack Compose 编写；分析层由 TextAnalyzer、SentenceImportanceAnalyzer、LocalReadingAgent 和 HistoryRepository 等对象承担；模型推理由 SentenceImportanceClassifier 封装；模型文件位于 `android/app/src/main/assets/sentence_importance_model.tflite`。

```mermaid
flowchart TD
    A[文本输入或示例文本] --> B[TextAnalyzer 本地摘要和关键词]
    B --> C[ArticleAnalysis 分析结果]
    C --> D[SentenceImportanceAnalyzer 提取5维特征]
    D --> E[SentenceImportanceClassifier LiteRT推理]
    C --> F[LocalReadingAgent 问答和复习材料]
    C --> G[HistoryRepository 本地历史记录]
    E --> H[Compose UI 展示]
    F --> H
    G --> H
```

### 3.2 Android 端架构

Android 工程使用 Gradle 管理，应用入口为 `MainActivity`。UI 通过 Compose 的状态驱动方式实现：用户输入文本后点击“开始智能分析”，界面状态更新为 ArticleAnalysis，摘要结果、LiteRT 结果和功能入口随状态显示。页面切换使用内部枚举 `Screen` 控制，包括 Home、AgentChat、KnowledgeCards 和 History。

### 3.3 模块划分

系统主要划分为七个模块：文本输入模块、本地摘要与关键词模块、LiteRT 端侧分析模块、Agent 问答模块、知识卡片与复习题模块、历史记录模块、UI 展示模块。模块之间通过 ArticleAnalysis 等数据结构传递结果，避免每个页面重复解析原文。

### 3.4 数据流设计

用户输入原文后，TextAnalyzer 先完成句子切分、关键词评分和摘要生成，得到 ArticleAnalysis。随后 SentenceImportanceAnalyzer 读取 ArticleAnalysis 中的原文和关键词，为文章前几句提取特征，并调用 SentenceImportanceClassifier 生成句子重要性分数。用户进入 Agent 或知识卡片页面时，LocalReadingAgent 基于同一份 ArticleAnalysis 生成回答、卡片和复习题。用户完成分析后，HistoryRepository 将结果保存到 SharedPreferences；从历史记录恢复时，再把历史原文交给分析流程重新生成当前页面状态。

### 3.5 版本迭代过程

项目按小版本逐步推进。v0.1 完成立项和仓库初始化；v0.2 做出摘要 MVP；v0.3 加入 Agent 问答和知识卡片；v0.4 完成 LiteRT 端侧模型集成；v0.5 加入本地历史记录；v1.0 冻结功能范围，整理 APK、截图、测试记录、演示材料和课程文档。这样的迭代方式使项目可以在每个阶段都有可运行成果，而不是最后一次性堆功能。

## 4 模块详细设计

### 4.1 文本输入模块

模块职责：提供待分析文本输入、示例文本选择、清空文本和开始分析入口。输入为用户粘贴的原文或内置 SampleTexts；输出为传递给 TextAnalyzer 的字符串。核心界面由 Compose 的 OutlinedTextField、Button 和状态变量组成。

处理流程：用户打开 App 后，首页显示文本输入框、示例按钮和分析按钮。点击示例文本时，输入框填入预设内容；点击开始智能分析时，系统校验文本是否为空，非空则进入分析流程，空文本则保留友好提示。异常处理重点是空文本、过短文本和现场误操作，不让 App 因输入问题崩溃。

### 4.2 本地摘要与关键词模块

模块职责：从原文中生成一句话总结、关键词和分点摘要。输入为原始文本；输出为 ArticleAnalysis 中的 oneSentenceSummary、keywords、bulletSummaries、sentenceCount 和 characterCount。核心对象为 `TextAnalyzer`。

处理流程：TextAnalyzer 使用正则切分句子，对领域词、英文技术词和中文短语进行评分，结合关键词命中、句子长度和句子位置给句子打分。得分最高的句子作为一句话总结，若干高分句子按原文顺序整理为分点摘要。异常处理包括文本为空返回 null、句子切分失败时截取原文前段作为兜底、关键词为空时提供默认关键词。

### 4.3 LiteRT 端侧句子重要性分析模块

模块职责：展示移动端机器学习能力，对文章中的句子给出重要性评分。输入为 ArticleAnalysis、关键词和句子位置；输出为 SentenceImportance 列表，包含句子、分数、等级、来源和索引。核心对象为 `SentenceImportanceAnalyzer` 和 `SentenceImportanceClassifier`。

处理流程：SentenceImportanceAnalyzer 对文章前 5 个句子提取 5 维特征：句子长度归一化、关键词重合度、位置分数、标点提示分数、总结提示词分数。SentenceImportanceClassifier 加载 assets 中的 `sentence_importance_model.tflite`，用 TensorFlow Lite Interpreter 推理，输出 0 到 1 的分数。分数大于等于 0.70 标为高，0.40 到 0.70 标为中，低于 0.40 标为低。

异常处理：如果模型文件不存在、Interpreter 初始化失败或推理异常，系统使用 fallbackScore 规则评分，并在结果来源中显示规则评分兜底。这样即使现场手机环境出现模型加载问题，摘要和演示流程仍能继续。

### 4.4 Agent 问答模块

模块职责：围绕当前文章分析结果回答用户问题。输入为用户问题和 ArticleAnalysis；输出为 Agent 回答文本。核心对象为 `LocalReadingAgent`，当前版本统一表述为“本地规则型阅读 Agent”。

处理流程：LocalReadingAgent 对问题进行意图识别，判断问题是否属于“文章主要讲什么”“关键词”“复习重点”“生成复习题”“知识卡片”等类型。不同意图调用 ArticleAnalysis 中的摘要、关键词、分点摘要或复习题生成逻辑组织回答。它不调用 GPT API，不依赖云端大模型，也不需要网络。

异常处理：如果用户尚未完成文章分析，Agent 会提示先输入文章并完成分析；如果用户提交空问题，会提示输入问题；如果问题不属于已知意图，则返回一个包含总结、关键词和分点摘要的综合回答。

### 4.5 知识卡片与复习题模块

模块职责：把分析结果转换成适合复习的卡片和题目。输入为 ArticleAnalysis；输出为 KnowledgeCard 列表和 QuizQuestion 列表。核心函数为 `generateKnowledgeCards` 和 `generateQuizQuestions`。

处理流程：知识卡片通常包括概念卡、核心观点卡和复习卡。复习题根据一句话总结、关键词和分点摘要生成简答题及参考答案。该模块的设计目的是把“读完文章得到摘要”推进到“能回看和复习”的阶段，提高 App 的完整度。

异常处理：如果没有分析结果，页面显示友好提示，不生成空卡片；如果分点摘要不足，则用一句话总结作为参考答案来源。

### 4.6 历史记录模块

模块职责：保存最近分析过的文章和结果，支持回看、恢复和清空。输入为 ArticleAnalysis；输出为 HistoryRecord 列表。核心对象为 `HistoryRepository`。

处理流程：每次分析完成后，系统把标题、预览、保存时间、原文、摘要、关键词、句子数和字符数写入 JSON，并存入 SharedPreferences。最多保留最近 12 条记录，重复原文会去重。用户点击历史记录中的恢复按钮时，系统取出原文并重新执行分析流程。

异常处理：如果 SharedPreferences 中的 JSON 解析失败，返回空列表；清空历史记录时写入空数组，避免残留脏数据影响演示。

### 4.7 UI 展示模块

模块职责：把输入区、结果区、LiteRT 区域、Agent 页面、知识卡片页面和历史页面组织成完整移动端体验。输入为各模块状态；输出为用户可操作界面。核心实现集中在 Compose 组件中。

处理流程：Home 页面负责输入和结果展示；AgentChat 页面负责快捷问题、手动提问和聊天记录；KnowledgeCards 页面负责展示卡片和复习题；History 页面负责展示最近记录、恢复和清空。UI 使用卡片式布局、分区标题和按钮入口，保证演示时老师能快速看到功能点。

## 5 核心算法与模型设计

### 5.1 本地摘要算法

本地摘要算法使用轻量规则实现，优点是稳定、可解释、无需网络。算法先切分句子，再统计关键词命中情况。句子得分由关键词命中、句子长度和位置共同决定。第一句通常有较高位置权重，中等长度句子更适合作为摘要候选。最终系统选择得分最高的句子作为一句话总结，再选择若干高分句子生成分点摘要。

这种方法不追求大型语言模型级别的自然语言生成效果，但非常适合课程项目演示。它能在没有网络和 API Key 的情况下稳定运行，也能让老师看到算法逻辑和代码实现，而不是只展示一个外部服务调用结果。

### 5.2 关键词提取算法

关键词提取结合三类线索：项目内置领域词、英文技术词和中文短语窗口。领域词包括 SmartRead、Agent、Android、Kotlin、Compose、LiteRT 等项目相关词，也包括摘要、关键词、知识卡片、问答、历史记录、模型等阅读分析词。英文技术词通过正则提取，中文短语通过双字窗口统计并过滤停用词。最终按分数排序，并避免包含关系过强的重复词。

### 5.3 LiteRT 模型设计

LiteRT 模型名称为 `sentence_importance_model`，输入形状为 `[1, 5]`，输出形状为 `[1, 1]`。模型特征为 `sentenceLengthNorm`、`keywordOverlapScore`、`positionScore`、`punctuationScore`、`summaryCueScore`。模型结构为 `Input(5) -> Dense(8, relu) -> Dense(4, relu) -> Dense(1, sigmoid)`，输出为 0.0 到 1.0 的句子重要性分数。

项目中的模型文件大小约 1496 字节，SHA256 为 `11462553e040257c11677bc7824021440761922ff4a8b3ced2bb7d21904b9dde`。训练摘要中记录的 weak label accuracy 为 0.9260，但这只是弱监督标签上的训练参考指标，演示时不把它夸大为真实业务准确率。该模型的作用是展示端侧模型集成和句子评分流程。

### 5.4 端侧推理流程

端侧推理流程如下：App 读取原文并生成 ArticleAnalysis；SentenceImportanceAnalyzer 取前 5 个句子；每个句子提取 5 维特征；SentenceImportanceClassifier 读取 assets 中的 .tflite 文件；Interpreter 执行推理；系统把输出分数转换为等级；界面展示句子、分数、等级和来源。模型加载失败时进入 fallbackScore，不中断主流程。

### 5.5 本地规则型 Agent 设计

本项目 Agent 当前版本没有调用 GPT API，也没有依赖云端大模型。它是本地规则型阅读 Agent，基于当前文章分析结果进行回答。它使用问题意图识别把用户输入映射到几个可解释动作：概括文章、提取关键词、整理复习重点、生成复习题、生成知识卡片和综合回答。这样设计的好处是离线可用、稳定、可控，也能避免课堂现场因为网络或 API 配置失败而影响演示。

## 6 数据结构设计

### ArticleAnalysis

ArticleAnalysis 是主分析结果结构，字段包括 originalText、oneSentenceSummary、keywords、bulletSummaries、sentenceCount、characterCount、sentenceImportances 和 localModelStatus。它把文本分析结果集中保存，供首页、Agent、知识卡片和历史记录模块复用。

### SentenceImportance

SentenceImportance 表示单个句子的端侧评分结果，字段包括 sentence、score、level、source 和 index。score 是模型或兜底规则输出的分数，level 是高、中、低，source 用于说明结果来自 LiteRT 本地模型还是规则评分兜底。

### ChatMessage

ChatMessage 用于 Agent 问答页面，字段包括 role 和 content。role 使用 MessageRole 区分 USER 和 AGENT。该结构让界面能够按用户消息和 Agent 回答分别展示气泡。

### KnowledgeCard

KnowledgeCard 用于知识卡片页面，字段包括 title、type 和 content。title 表示卡片标题，type 表示概念、观点或复习等类别，content 是卡片正文。

### QuizQuestion

QuizQuestion 用于复习题展示，字段包括 question 和 referenceAnswer。当前版本以简答题为主，参考答案来自摘要、关键词和分点摘要。

### HistoryRecord

HistoryRecord 用于历史记录，字段包括 id、title、preview、savedAt、displayTime、originalText、oneSentenceSummary、keywords、sentenceCount 和 characterCount。它兼顾列表展示和恢复分析所需数据。

## 7 界面设计

### 7.1 首页

首页展示项目名称、版本说明、文本输入区、示例文本按钮、开始分析按钮、清空按钮、字符数和句子数。截图参考：`screenshots/app/V1.0首页_20260512.png`。

### 7.2 文本输入区

文本输入区支持粘贴或选择示例文本，提示建议输入 100 字以上。演示时为了稳定，优先使用示例文本。截图参考：`screenshots/app/V1.0文本输入_20260512.png`。

### 7.3 摘要结果区

摘要结果区展示一句话总结、关键词、分点摘要和入口按钮。截图参考：`screenshots/app/V1.0摘要结果_20260512.png`。

### 7.4 LiteRT 分析区

LiteRT 区域展示推理方式、句子重要性等级、分数和来源，是演示中机器学习整合的重点。截图参考：`screenshots/app/V1.0LiteRT端侧分析_20260512.png`。

### 7.5 Agent 问答页

Agent 页面展示当前文章摘要、快捷问题、聊天记录和手动提问输入框。截图参考：`screenshots/app/V1.0Agent问答_20260512.png`、`screenshots/app/V1.0Agent快捷问题_20260512.png`。

### 7.6 知识卡片页

知识卡片页面展示概念卡、核心观点卡、复习卡和复习题。截图参考：`screenshots/app/V1.0知识卡片_20260512.png`、`screenshots/app/V1.0复习题_20260512.png`。

### 7.7 历史记录页

历史记录页展示最近保存的分析，支持恢复历史分析和清空历史记录。截图参考：`screenshots/app/V1.0历史记录页_20260512.png`、`screenshots/app/V1.0历史记录恢复_20260512.png`。

## 8 测试设计

### 8.1 测试环境

开发环境为 Windows + Android Studio / Gradle，Android 工程使用 Kotlin 和 Jetpack Compose。模拟器测试使用 Pixel_8 Android 模拟器；真机测试使用安卓手机安装 V1.0 debug APK。APK 文件位于 `release/SmartReadAgent-v1.0-debug.apk`。

### 8.2 测试用例

测试覆盖 App 启动、示例文本输入、手动文本输入、开始智能分析、空文本提示、摘要结果展示、关键词展示、LiteRT 区域展示、Agent 快捷问题、Agent 手动提问、知识卡片、复习题、历史记录保存、历史记录恢复和清空历史记录。

### 8.3 测试结果

V1.0 阶段模拟器测试和真机体验反馈均表明主要演示流程可用。示例文本能够生成摘要，LiteRT 分析区域能够显示，Agent 问答、知识卡片、复习题和历史记录均可进入。当前没有记录影响现场演示的阻塞性 Bug。

### 8.4 Bug 与优化建议

当前主要优化建议包括：摘要算法仍偏轻量，后续可引入更强的抽取式或生成式方法；LiteRT 模型训练数据规模有限，后续可扩充真实标注数据；历史记录当前使用 SharedPreferences，后续可迁移到 Room；正式发布时需要生成正式签名 release 包；如果需要处理图片资料，可增加 OCR 导入。

## 9 部署与运行

### 9.1 Android Studio 构建

Android 工程位于 `android/`。构建命令为：

```powershell
cd android
.\\gradlew.bat :app:assembleDebug
```

### 9.2 APK 安装

课程演示 APK 位于 `release/SmartReadAgent-v1.0-debug.apk`。安卓手机安装时需要允许当前应用安装未知来源应用。该 APK 是课程演示用 debug APK，不作为应用商店正式商业发布版。

### 9.3 GitHub/Gitee 仓库

项目使用 Gitee 和 GitHub 进行版本管理，保留了 main、dev、feature 分支和 v0.1 到 v1.0 的阶段记录。仓库用于展示源码、提交记录、标签和 Release 材料。

### 9.4 演示环境

期末演示时间为 2026-07-01 9:50-10:05，地点为计网大楼 202。建议 12 分钟完成 PPT 和 App 演示，预留 3 分钟回答问题。林立洲负责 PPT 汇报和技术说明，江轩宇负责安卓手机实际操作演示。现场准备手机、APK、PPT、模拟器备份、截图备份和答辩问题。

## 10 项目总结与后续优化

SmartRead Agent V1.0 已经完成课程期末演示所需的核心能力：它能运行在 Android 端，能通过示例文本稳定展示摘要、关键词、分点摘要、LiteRT 端侧分析、Agent 问答、知识卡片、复习题和历史记录。项目最重要的特点是把移动端应用和端侧机器学习结合起来，而不是只做普通文本处理页面。LiteRT 模型虽然轻量，但已经走通了特征提取、模型导出、Android assets 集成、本地推理和 UI 展示的完整流程。

从项目分工看，林立洲主要承担选题、技术方案、Android 核心开发、LiteRT 集成、版本管理、文档和汇报材料整理；江轩宇主要承担需求反馈、真机安装体验、功能流程验证、UI 使用建议、文档校对和现场演示操作。这样的分工符合双人小组的实际情况，也能在演示中清楚说明各自贡献。

后续如果继续完善，可以从四个方向推进。第一，扩充 LiteRT 模型训练数据，用更真实的阅读材料和人工标注提升句子重要性判断效果。第二，完善数据存储，把 SharedPreferences 迁移到 Room 数据库，支持更完整的历史记录管理。第三，加入 OCR 或文件导入，让用户能从图片、PDF 或课堂讲义中导入文本。第四，在保持本地能力稳定的基础上，接入云端大模型作为增强问答能力，但必须明确这是后续增强，不是当前 V1.0 已实现内容。
"""


def ensure_dirs() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)
    DEMO.mkdir(parents=True, exist_ok=True)
    PM_MEMO.parent.mkdir(parents=True, exist_ok=True)


def write(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.strip() + "\n", encoding="utf-8")


def build_ppt_markdown() -> str:
    parts = [
        "# SmartRead Agent 期末演示 PPT 逐页内容",
        "",
        "- 课程：软件项目研发实践（3）",
        "- 演示顺序：第4组，2026-07-01 9:50-10:05，计网大楼 202",
        "- 演示控制：12 分钟讲解与操作，3 分钟答问",
        "- 汇报分工：林立洲负责 PPT 汇报和技术说明；江轩宇负责手机实际操作演示",
        "- 统一口径：Agent 是本地规则型阅读 Agent，当前版本没有调用 GPT API，也没有依赖云端大模型",
    ]
    for index, slide in enumerate(SLIDES, 1):
        parts.extend(
            [
                "",
                f"## {index}. {slide['title']}",
                "",
                f"副标题：{slide['subtitle']}",
                "",
                "要点：",
            ]
        )
        parts.extend([f"- {item}" for item in slide["bullets"]])
        if slide.get("image"):
            parts.append(f"- 建议配图：`screenshots/app/{slide['image']}`")
        if slide.get("diagram"):
            parts.append("- 建议配图：系统架构流程图（文本输入 -> 摘要 -> LiteRT -> Agent -> 历史记录）")
        parts.extend(
            [
                "",
                "讲解提示：用 30-60 秒讲清本页，不要照着大段文字念。每页围绕评分表中的一个得分点展开。",
            ]
        )
    return "\n".join(parts)


def build_talk_script() -> str:
    return """# SmartRead Agent 期末现场演示讲稿

## 一、演示基本信息

- 课程：软件项目研发实践（3）
- 演示顺序：第4组 SmartRead Agent
- 时间地点：2026-07-01 9:50-10:05，计网大楼 202
- 时间控制：建议 12 分钟完成展示，预留 3 分钟回答问题
- 汇报分工：林立洲负责 PPT 汇报、技术说明和答辩主答；江轩宇负责安卓手机实际操作演示
- 统一口径：本项目当前 Agent 是本地规则型阅读 Agent，不调用 GPT API，不依赖云端大模型

## 二、12 分钟时间安排

| 时间 | 内容 | 主讲/操作 |
|---|---|---|
| 0:00-0:40 | 封面、成员分工、演示安排 | 林立洲 |
| 0:40-1:40 | 项目背景、痛点和目标 | 林立洲 |
| 1:40-2:40 | 功能总览 | 林立洲 |
| 2:40-4:20 | 技术架构、LiteRT 模型设计 | 林立洲 |
| 4:20-5:10 | Agent 实现说明 | 林立洲 |
| 5:10-8:30 | App 实际操作演示 | 江轩宇操作，林立洲说明 |
| 8:30-9:40 | 迭代过程、测试结果 | 林立洲 |
| 9:40-10:50 | 团队协作和项目管理证据 | 林立洲 |
| 10:50-11:40 | 项目亮点、不足和后续优化 | 林立洲 |
| 11:40-12:00 | 收尾，进入问答 | 林立洲 |

## 三、逐页讲解稿

### 1. 封面

大家好，我们是第4组，项目名称是 SmartRead Agent：AI 智读助手。这个项目是一个 Android/Kotlin 移动端阅读辅助应用，主要面向学生阅读教材、课堂资料和技术文章的场景。我是林立洲，负责汇报、核心开发、LiteRT 集成和材料整理；江轩宇负责手机端实际操作、真机测试、使用反馈和演示配合。

### 2. 项目背景与痛点

我们选择这个方向，是因为学习过程中经常会遇到较长的资料。手工整理摘要、关键词和复习题比较耗时，而且容易遗漏重点。SmartRead Agent 想解决的不是“替代学习”，而是帮助用户更快完成第一遍梳理。

### 3. 项目目标

项目目标可以概括为一个闭环：输入文本，生成摘要，做端侧句子重要性分析，再用 Agent 问答和知识卡片辅助复习，最后通过历史记录回看。这个闭环能比较完整地展示一个移动 AI 阅读助手的基本形态。

### 4. 功能总览

当前 V1.0 支持示例文本、手动输入、一句话总结、关键词、分点摘要、LiteRT 分析、Agent 问答、知识卡片、复习题和历史记录。演示时我们会按这个顺序走一遍，确保老师能看到所有核心功能。

### 5. 技术架构

技术上，App 使用 Android/Kotlin 和 Jetpack Compose。文本处理由 TextAnalyzer 完成，Agent 由 LocalReadingAgent 完成，LiteRT 推理由 SentenceImportanceClassifier 完成，历史记录由 SharedPreferences 保存。整个 V1.0 版本的核心流程都在本地完成，不依赖网络。

### 6. 机器学习与 LiteRT 设计

评分表中机器学习整合占 40%，所以这里是重点。我们的 LiteRT 模型做的是句子重要性分析。系统从句子中提取 5 维特征，包括长度、关键词重合、位置、标点和提示词，然后输入 .tflite 模型，输出一个 0 到 1 的重要性分数。模型文件随 APK 放在 assets 中，手机端本地推理。

### 7. Agent 问答实现说明

这里需要特别说明：当前版本没有调用 GPT API，也没有依赖云端大模型。Agent 是本地规则型阅读 Agent。它根据当前文章的摘要、关键词、分点摘要和复习题规则，通过问题意图识别来回答。优点是离线可用、稳定、可解释，适合现场演示。

### 8. App 演示流程

接下来由江轩宇操作手机，我同步说明。我们会先打开 App，点击示例文本，再点击开始智能分析。然后展示摘要、关键词、分点摘要和 LiteRT 区域。接着进入 Agent 问答页，点击快捷问题，再进入知识卡片页看卡片和复习题，最后进入历史记录页恢复一次历史分析。

### 9. 项目迭代过程

项目从 v0.1 到 v1.0 是逐步推进的。v0.2 先完成摘要 MVP，v0.3 加 Agent 和知识卡片，v0.4 加 LiteRT，v0.5 加历史记录，v1.0 收口 APK、截图、测试和演示材料。这样每个阶段都有可运行成果，也方便版本管理。

### 10. 团队协作与分工

我们是双人小组，分工按实际参与情况说明。林立洲主要负责核心开发、LiteRT、Agent、版本管理和汇报材料；江轩宇主要负责需求反馈、真机安装、功能体验、UI 建议、文档校对和演示操作。过程材料中也保留了 QQ 沟通、项目管理平台和 Gitee/GitHub 记录。

### 11. 测试与运行结果

测试方面，我们做了模拟器测试和安卓真机体验测试。覆盖 App 启动、示例文本分析、LiteRT 展示、Agent 问答、知识卡片、复习题和历史记录。当前没有影响演示的阻塞性 Bug。

### 12. 项目亮点

项目亮点有三个：第一，应用场景贴近学生学习；第二，功能闭环比较完整，不只是单个摘要页面；第三，LiteRT 模型真正集成到 Android 端侧运行，能够体现移动机器学习应用。

### 13. 不足与后续优化

不足也比较明确。摘要和 Agent 当前比较轻量，模型数据规模有限，历史记录还没有使用数据库。后续可以提升模型训练数据，引入 Room 数据库，增加 OCR，或者在本地能力稳定的基础上接入云端大模型增强问答。

### 14. 结束页

以上就是我们的项目展示。SmartRead Agent V1.0 是课程演示稳定版，APK、源码、截图、PPT 和详细设计说明书都已经整理。谢谢老师和同学，欢迎提问。

## 四、江轩宇手机操作步骤

1. 提前确认手机已安装 `SmartReadAgent-v1.0-debug.apk`。
2. 打开 SmartRead Agent，停留 2 秒，让老师看到首页和版本号。
3. 点击“示例 1”或提前约定的示例文本。
4. 点击“开始智能分析”。
5. 慢慢向下滑，依次展示一句话总结、关键词、分点摘要。
6. 停在 LiteRT 端侧分析区域，等待林立洲说明 5 维特征和端侧推理。
7. 点击进入 Agent 问答页。
8. 点击“这篇文章主要讲了什么”或“帮我整理复习重点”等快捷问题。
9. 返回摘要页，进入知识卡片页，展示卡片和复习题。
10. 返回摘要页，进入历史记录页，点击恢复一条历史分析。
11. 回到首页，准备答辩。

## 五、问答时优先守住的口径

- Agent：本地规则型阅读 Agent，不调用 GPT API。
- LiteRT：端侧句子重要性分析模型，assets 集成，手机本地推理。
- 分工：林立洲负责核心开发和汇报，江轩宇负责测试体验和现场操作，不夸大核心算法贡献。
- APK：课程演示用 debug APK，不是应用商店商业发布版。
"""


def build_qa_markdown() -> str:
    lines = [
        "# SmartRead Agent 答辩可能问题",
        "",
        "统一口径：当前版本的 Agent 是本地规则型阅读 Agent，没有调用 GPT API，也没有依赖云端大模型。LiteRT 模块是本项目移动机器学习整合的重点。",
    ]
    for idx, (question, answer) in enumerate(QA, 1):
        lines.extend([f"", f"## {idx}. {question}", "", answer])
    return "\n".join(lines)


def build_score_check() -> str:
    return """# 期末评分项对应自检表

| 评分项 | 分值占比 | 项目对应内容 | 演示中如何体现 | 材料证据 |
|---|---:|---|---|---|
| APP 应用创新和实用性 | 20% | 面向学生阅读教材、课堂资料和技术文章，解决摘要整理、重点提取和复习回看的问题 | 第2-3页讲背景和目标，App 演示使用示例文本展示完整阅读辅助流程 | README、项目说明文档、首页和摘要截图 |
| 移动应用功能丰富度 | 20% | 文本输入、示例文本、摘要、关键词、分点摘要、Agent、知识卡片、复习题、历史记录 | 第4页列功能总览，第8页现场依次操作摘要、Agent、卡片、历史记录 | V1.0 各页面截图、APK、演示讲稿 |
| 机器学习在应用中的整合程度 | 40% | LiteRT 端侧句子重要性分析，5 维特征，.tflite 模型随 APK 集成，本地推理并显示评分 | 第6页重点讲 LiteRT，演示时停留在 LiteRT 区域解释分数和来源 | `model/exports/model_metadata.json`、`android/app/src/main/assets/sentence_importance_model.tflite`、LiteRT 截图 |
| 团队成员项目协作度和精神面貌 | 10% | 林立洲负责核心开发和汇报，江轩宇负责真机体验、功能验证、UI 建议、文档校对和现场操作 | 第10页讲真实分工，现场由江轩宇配合操作手机 | QQ 沟通记录、项目管理平台截图、Gitee/GitHub 版本记录 |
| PPT 制作、听众主题效果及时间控制 | 10% | 15页 PPT，12分钟讲解和操作，3分钟答问，文字简洁并配截图 | 按讲稿控制节奏，重点页为 LiteRT、Agent 和运行效果展示 | `SmartReadAgent_期末演示PPT.pptx`、`期末现场演示讲稿.md` |

## 演示优先级

1. 一定要讲清楚 LiteRT，因为它占机器学习整合 40%。
2. 一定要说清楚 Agent 是本地规则型阅读 Agent，不调用 GPT API。
3. 一定要让老师看到真实 App，而不是只讲 PPT。
4. 分工要自然说明，不写贡献比例在 PPT 上；比例等演示结束后按老师要求提交给班长。
5. 如果现场时间紧，宁愿少讲项目管理细节，也要完整展示 LiteRT 和主要功能闭环。
"""


def build_day_checklist() -> str:
    return """# 演示当天检查清单

## 一、时间地点

- [ ] 2026-07-01 上午提前 5 分钟到场。
- [ ] 第4组 SmartRead Agent，预计 9:50-10:05。
- [ ] 地点：计网大楼 202。
- [ ] 服仪整齐，不穿拖鞋入场。

## 二、文件准备

- [ ] PPT：`docs/SmartReadAgent_期末演示PPT.pptx`。
- [ ] APK：`release/SmartReadAgent-v1.0-debug.apk`。
- [ ] 详细设计说明书初稿：`docs/SmartReadAgent_软件详细设计说明书_初稿.docx`。
- [ ] 演示讲稿：`docs/期末现场演示讲稿.md`。
- [ ] 答辩问题：`docs/答辩可能问题.md`。
- [ ] 评分项自检表：`docs/期末评分项对应自检表.md`。
- [ ] App 截图备份：`screenshots/app/`。

## 三、设备准备

- [ ] 笔记本电脑和充电器。
- [ ] 安卓手机已安装 `SmartReadAgent-v1.0-debug.apk`。
- [ ] 手机电量充足，亮度调高，关闭无关通知。
- [ ] 模拟器或截图备份可用。
- [ ] PPT 能正常打开，图片不丢失。
- [ ] 如果投屏，提前测试横竖屏和窗口大小。

## 四、演示流程确认

- [ ] 林立洲负责 PPT 汇报和答辩主答。
- [ ] 江轩宇负责手机操作演示。
- [ ] 操作顺序：打开 App -> 示例文本 -> 开始智能分析 -> 摘要 -> LiteRT -> Agent -> 知识卡片 -> 历史记录。
- [ ] LiteRT 页面至少停留 20 秒，讲清楚端侧模型和 5 维特征。
- [ ] Agent 页面说清楚“本地规则型阅读 Agent，不调用 GPT API”。
- [ ] 演示控制在 12 分钟左右，预留 3 分钟问答。

## 五、演示后提交

- [ ] 演示结束后按老师要求向班长提交小组成员得分比例，总和为人数 x 100%。
- [ ] 2026-07-10 前提交详细设计说明书电子版和纸质版。
- [ ] 保留学习通或群文件提交截图。
"""


def build_final_submit() -> str:
    return """# SmartRead Agent 最终提交说明

## 一、课程演示要求

根据《软件项目研发实践（3）》期末考核说明，本组为第4组，演示时间为 2026-07-01 9:50-10:05，地点为计网大楼 202。每组总时长不超过 15 分钟，建议 12 分钟完成 PPT 汇报和系统演示，预留 3 分钟回答问题。

演示需要简要介绍作品功能和小组成员分工，然后进行系统演示。PPT 无固定模板，但必须展示作品功能，并包含小组分工。

## 二、演示材料

- PPT：`docs/SmartReadAgent_期末演示PPT.pptx`
- PPT 逐页内容：`docs/期末演示PPT逐页内容.md`
- 现场演示讲稿：`docs/期末现场演示讲稿.md`
- 答辩问题准备：`docs/答辩可能问题.md`
- 评分项自检表：`docs/期末评分项对应自检表.md`
- 演示当天检查清单：`docs/演示当天检查清单.md`
- App 截图：`screenshots/app/`

## 三、应用与源码材料

- Android 工程：`android/`
- 模型材料：`model/`
- V1.0 APK：`release/SmartReadAgent-v1.0-debug.apk`
- APK 信息说明：`release/APK信息说明.md`
- Android 手机安装说明：`release/Android手机安装说明.md`
- v1.0 Release 说明：`release/v1.0_release_notes.md`

说明：`SmartReadAgent-v1.0-debug.apk` 是课程演示用 debug APK，用于课程验收、真机测试和现场演示，不作为应用商店正式商业发布版本。

## 四、详细设计说明书

演示结束后还需要提交移动机器学习应用的详细设计说明书。电子版和纸质版需在 2026-07-10 前提交给班长。

- Markdown 初稿：`docs/SmartReadAgent_软件详细设计说明书_初稿.md`
- Word 初稿：`docs/SmartReadAgent_软件详细设计说明书_初稿.docx`

详细设计说明书已覆盖项目背景、总体架构、模块设计、LiteRT 端侧模型、Agent 实现、数据结构、界面设计、测试设计、部署运行和后续优化。

## 五、小组分工说明

- 林立洲：项目负责人，负责整体选题与技术方案、Android 核心功能开发、本地摘要、关键词、Agent 问答、LiteRT 模型集成、GitHub/Gitee 版本管理、项目文档、PPT 汇报和最终材料整理。
- 江轩宇：负责需求反馈与使用体验测试、安卓真机安装与功能验证、Agent 问答/知识卡片/历史记录等功能的演示操作、UI 使用建议、文档校对和现场演示配合。

PPT 中只说明真实分工，不写 105% / 95% 贡献比例；演示结束后按老师要求另行向班长提交小组成员得分比例。

## 六、重点答辩口径

- Agent：当前版本是本地规则型阅读 Agent，没有调用 GPT API，也没有依赖云端大模型。
- LiteRT：模型用于端侧句子重要性分析，输入 5 维特征，`.tflite` 文件集成到 Android assets，本地推理。
- 离线能力：摘要、关键词、Agent、知识卡片、复习题、历史记录和 LiteRT 分析均可离线运行。
- 未实现内容：当前版本没有登录、服务器、OCR、云端同步和正式商店发布，这些是后续优化方向。
"""


def build_pm_memo() -> str:
    return """# 个人项目管理报告素材备忘

## 期末演示评分表导向

《软件项目研发实践（3）》期末评分表中，APP 应用创新和实用性占 20%，移动应用功能丰富度占 20%，机器学习在应用中的整合程度占 40%，团队成员项目协作度和精神面貌占 10%，PPT 制作、听众主题效果及时间控制占 10%。因此项目管理报告和后续复盘可以重点写：项目为什么贴近学生阅读场景、功能如何从 MVP 逐步扩展、LiteRT 如何体现移动机器学习、两人如何协作完成演示和测试。

## 小组分工

- 林立洲：项目负责人，负责整体选题、技术方案、Android/Kotlin 核心开发、本地摘要、关键词、Agent 问答、LiteRT 模型集成、GitHub/Gitee 版本管理、项目文档、PPT 汇报和最终材料整理。
- 江轩宇：负责需求反馈、安卓真机安装与功能验证、Agent 问答、知识卡片、历史记录等功能的演示操作、UI 使用建议、文档校对和现场演示配合。

## 江轩宇参与内容

江轩宇的参与重点不是核心算法开发，而是最终演示与测试验收模块。他参与了真机安装、功能流程体验、Agent 问答和知识卡片入口确认、历史记录恢复流程体验、UI 使用建议、文档校对和演示操作配合。报告中可以写他对“能不能装、能不能顺利演示、页面流程是否容易讲清楚”这类问题提供反馈。

## 真机测试

V1.0 阶段生成 `SmartReadAgent-v1.0-debug.apk` 后，进行了安卓真机体验反馈。测试重点包括安装启动、示例文本摘要、LiteRT 端侧分析展示、Agent 问答、知识卡片、复习题、历史记录保存和恢复。当前没有影响演示的阻塞性 Bug。

## 演示配合

软件实践 3 的现场演示安排在 2026-07-01，SmartRead Agent 为第4组，时间 9:50-10:05。演示分工为林立洲负责 PPT 汇报和技术说明，江轩宇负责手机端实际操作。演示流程固定为：首页 -> 示例文本 -> 开始智能分析 -> 摘要 -> LiteRT -> Agent -> 知识卡片 -> 历史记录。

## 项目管理平台

项目管理平台用于记录阶段任务、任务拆分、任务状态、Bug 管理和进度复盘。报告中可以结合 Leangoo 看板截图说明，小组不是只在最后提交代码，而是在不同阶段把功能、测试和材料整理拆成任务推进。

## Gitee/GitHub 迭代过程

项目保留了从 v0.1 到 v1.0 的版本推进。v0.1 完成立项和仓库初始化；v0.2 完成摘要 MVP；v0.3 完成 Agent 问答和知识卡片；v0.4 完成 LiteRT 端侧模型集成；v0.5 完成历史记录；v1.0 完成演示稳定版、APK、截图、测试和文档收口。Gitee/GitHub 提交记录、标签和分支截图可以作为版本管理证据。

## 报告可写的真实反思

1. 双人小组任务集中，后续应更早拆分文档、测试和演示任务，避免最后阶段集中处理。
2. 端侧机器学习部分应更早确定模型输入和展示方式，否则容易出现“模型有了但讲不清楚”的问题。
3. 真机安装测试很重要，模拟器能跑不代表课堂演示一定稳定。
4. 项目管理材料最好随版本同步整理，不要等到最后一次性补截图。
5. 分工说明要按真实参与情况写，既不能把测试和演示配合写没了，也不能把非核心开发成员写成核心算法负责人。
"""


def build_markdown_files() -> None:
    write(DOCS / "期末演示PPT逐页内容.md", build_ppt_markdown())
    talk = build_talk_script()
    write(DOCS / "期末现场演示讲稿.md", talk)
    write(DOCS / "期末演示讲稿.md", talk)
    write(DOCS / "答辩可能问题.md", build_qa_markdown())
    write(DOCS / "SmartReadAgent_软件详细设计说明书_初稿.md", DESIGN_MD)
    write(DOCS / "期末评分项对应自检表.md", build_score_check())
    write(DOCS / "演示当天检查清单.md", build_day_checklist())
    write(DOCS / "最终提交说明.md", build_final_submit())
    write(PM_MEMO, build_pm_memo())


def add_textbox(slide, left, top, width, height, text, size=20, bold=False, color=(31, 41, 55)):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = RGBColor(*color)
    return shape


def add_bullets(slide, bullets, left, top, width, height):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(18)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = RGBColor(55, 65, 81)
        p.space_after = PptPt(8)
    return shape


def add_picture_fit(slide, image_path: Path, left, top, max_width, max_height):
    if not image_path.exists():
        return None
    pic = slide.shapes.add_picture(str(image_path), left, top, height=max_height)
    if pic.width > max_width:
        ratio = max_width / pic.width
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = int(left + (max_width - pic.width) / 2)
    pic.top = int(top + (max_height - pic.height) / 2)
    return pic


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    bg = RGBColor(248, 250, 252)
    blue = RGBColor(29, 78, 216)

    for idx, data in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
        add_textbox(slide, PptInches(0.55), PptInches(0.35), PptInches(12.2), PptInches(0.45), f"{idx:02d}", 13, True, (37, 99, 235))
        add_textbox(slide, PptInches(0.8), PptInches(0.7), PptInches(11.8), PptInches(0.55), data["title"], 30, True, (17, 24, 39))
        add_textbox(slide, PptInches(0.82), PptInches(1.22), PptInches(11.2), PptInches(0.35), data["subtitle"], 15, False, (75, 85, 99))
        if data.get("image"):
            add_bullets(slide, data["bullets"], PptInches(0.9), PptInches(1.85), PptInches(6.2), PptInches(4.9))
            add_picture_fit(slide, SCREENSHOTS / data["image"], PptInches(7.6), PptInches(1.55), PptInches(4.7), PptInches(5.45))
        elif data.get("diagram"):
            add_bullets(slide, data["bullets"], PptInches(0.9), PptInches(1.75), PptInches(5.8), PptInches(5.0))
            labels = ["文本输入", "TextAnalyzer", "LiteRT模型", "本地Agent", "历史记录", "Compose界面"]
            x_positions = [7.0, 9.05, 11.1, 8.0, 10.05, 9.05]
            y_positions = [1.9, 1.9, 1.9, 3.55, 3.55, 5.15]
            for label, x, y in zip(labels, x_positions, y_positions):
                shape = slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(1.6), PptInches(0.55))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(219, 234, 254)
                shape.line.color.rgb = blue
                p = shape.text_frame.paragraphs[0]
                p.text = label
                p.alignment = PP_ALIGN.CENTER
                p.font.size = PptPt(14)
                p.font.name = "Microsoft YaHei"
                p.font.bold = True
                p.font.color.rgb = RGBColor(30, 64, 175)
            add_textbox(slide, PptInches(7.2), PptInches(6.0), PptInches(5.4), PptInches(0.5), "核心链路：输入 -> 分析 -> 端侧模型 -> 问答复习 -> 保存回看", 15, False, (75, 85, 99))
        else:
            add_bullets(slide, data["bullets"], PptInches(1.0), PptInches(1.9), PptInches(11.2), PptInches(4.8))
        add_textbox(slide, PptInches(0.8), PptInches(6.95), PptInches(8.5), PptInches(0.25), "SmartRead Agent V1.0 期末演示稳定版", 10, False, (107, 114, 128))
    prs.save(DOCS / "SmartReadAgent_期末演示PPT.pptx")


def style_doc(document: Document) -> None:
    styles = document.styles
    styles["Normal"].font.name = "宋体"
    styles["Normal"].font.size = Pt(10.5)
    for name in ["Heading 1", "Heading 2", "Heading 3"]:
        styles[name].font.name = "黑体"


def add_markdown_to_doc(document: Document, markdown: str) -> None:
    table_buffer = []

    def flush_table():
        nonlocal table_buffer
        if not table_buffer:
            return
        rows = []
        for line in table_buffer:
            cells = [c.strip() for c in line.strip("|").split("|")]
            if all(set(c) <= {"-", ":"} for c in cells):
                continue
            rows.append(cells)
        if rows:
            table = document.add_table(rows=len(rows), cols=max(len(r) for r in rows))
            table.style = "Table Grid"
            for r_idx, row in enumerate(rows):
                for c_idx, value in enumerate(row):
                    table.cell(r_idx, c_idx).text = value
        table_buffer = []

    for raw in markdown.splitlines():
        line = raw.rstrip()
        if line.startswith("|") and line.endswith("|"):
            table_buffer.append(line)
            continue
        flush_table()
        if not line:
            continue
        if line.startswith("# "):
            document.add_heading(line[2:].strip(), level=0)
        elif line.startswith("## "):
            document.add_heading(line[3:].strip(), level=1)
        elif line.startswith("### "):
            document.add_heading(line[4:].strip(), level=2)
        elif line.startswith("#### "):
            document.add_heading(line[5:].strip(), level=3)
        elif line.startswith("- "):
            document.add_paragraph(line[2:].strip(), style="List Bullet")
        elif re.match(r"^\d+\. ", line):
            document.add_paragraph(re.sub(r"^\d+\. ", "", line), style="List Number")
        elif line.startswith("```"):
            continue
        else:
            document.add_paragraph(line)
    flush_table()


def build_docx() -> None:
    doc = Document()
    style_doc(doc)
    add_markdown_to_doc(doc, DESIGN_MD)
    doc.add_page_break()
    doc.add_heading("附录：界面截图", level=1)
    screenshots = [
        ("首页", "V1.0首页_20260512.png"),
        ("摘要结果", "V1.0摘要结果_20260512.png"),
        ("LiteRT 端侧分析", "V1.0LiteRT端侧分析_20260512.png"),
        ("Agent 问答", "V1.0Agent问答_20260512.png"),
        ("知识卡片", "V1.0知识卡片_20260512.png"),
        ("历史记录", "V1.0历史记录页_20260512.png"),
    ]
    for title, filename in screenshots:
        path = SCREENSHOTS / filename
        doc.add_heading(title, level=2)
        p = doc.add_paragraph(f"截图文件：screenshots/app/{filename}")
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if path.exists():
            doc.add_picture(str(path), width=Inches(3.2))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    try:
        doc.save(DOCS / "SmartReadAgent_软件详细设计说明书_初稿.docx")
    except PermissionError:
        doc.save(DOCS / "SmartReadAgent_软件详细设计说明书_自动生成临时版.docx")


def main() -> None:
    ensure_dirs()
    build_markdown_files()
    build_pptx()
    build_docx()
    # The teacher template has a fixed detailed-design structure. Rebuild that
    # document last so this broad generator cannot overwrite it with the looser
    # narrative version.
    template_rebuilder = PROJECT / "tools" / "rebuild_template_design_doc.py"
    if template_rebuilder.exists():
        subprocess.run([sys.executable, str(template_rebuilder)], check=True)
    formal_ppt_rebuilder = PROJECT / "tools" / "rebuild_formal_ppt.py"
    if formal_ppt_rebuilder.exists():
        subprocess.run([sys.executable, str(formal_ppt_rebuilder)], check=True)
    print("generated:")
    for path in [
        DOCS / "期末演示PPT逐页内容.md",
        DOCS / "SmartReadAgent_期末演示PPT.pptx",
        DOCS / "期末现场演示讲稿.md",
        DOCS / "答辩可能问题.md",
        DOCS / "SmartReadAgent_软件详细设计说明书_初稿.md",
        DOCS / "SmartReadAgent_软件详细设计说明书_初稿.docx",
        DOCS / "期末评分项对应自检表.md",
        DOCS / "演示当天检查清单.md",
        DOCS / "最终提交说明.md",
        PM_MEMO,
    ]:
        print(path)


if __name__ == "__main__":
    main()
